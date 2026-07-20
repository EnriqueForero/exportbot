"""Presentación institucional de un resultado de ExportBot (python-pptx)."""

from __future__ import annotations

import io
from datetime import datetime
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from config import VERSION_APP
from exportadores import ADVERTENCIA_LEGAL, AMARILLO_ACENTO, AZUL_INSTITUCIONAL

_AZUL = RGBColor.from_string(AZUL_INSTITUCIONAL)
_AMARILLO = RGBColor.from_string(AMARILLO_ACENTO)
_FILAS_POR_LAMINA = 10
_MAX_LAMINAS_TABLA = 3
_MAX_COLUMNAS = 8


def _lamina_titulo(prs: Presentation, pregunta: str, texto: str) -> None:
    lam = prs.slides.add_slide(prs.slide_layouts[6])
    banda = lam.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.1))
    banda.fill.solid()
    banda.fill.fore_color.rgb = _AZUL
    banda.line.fill.background()
    tf = banda.text_frame
    tf.text = "ExportBot 2.0 · ProColombia"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    cuadro = lam.shapes.add_textbox(Inches(0.6), Inches(1.5), prs.slide_width - Inches(1.2), Inches(4.6))
    t = cuadro.text_frame
    t.word_wrap = True
    p1 = t.paragraphs[0]
    p1.text = f"Pregunta: {pregunta}"
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = _AZUL
    p2 = t.add_paragraph()
    p2.text = texto
    p2.font.size = Pt(14)
    p3 = t.add_paragraph()
    p3.text = f"Generado: {datetime.now():%Y-%m-%d %H:%M} · Versión: {VERSION_APP}"
    p3.font.size = Pt(10)


def _laminas_tabla(prs: Presentation, columnas: list[str], filas: list[list[Any]]) -> None:
    columnas = columnas[:_MAX_COLUMNAS]
    for bloque in range(0, min(len(filas), _FILAS_POR_LAMINA * _MAX_LAMINAS_TABLA), _FILAS_POR_LAMINA):
        trozo = filas[bloque : bloque + _FILAS_POR_LAMINA]
        lam = prs.slides.add_slide(prs.slide_layouts[6])
        forma = lam.shapes.add_table(
            rows=len(trozo) + 1,
            cols=len(columnas),
            left=Inches(0.4),
            top=Inches(0.5),
            width=prs.slide_width - Inches(0.8),
            height=Inches(0.4) * (len(trozo) + 1),
        )
        tabla = forma.table
        for j, col in enumerate(columnas):
            celda = tabla.cell(0, j)
            celda.text = str(col)
            celda.fill.solid()
            celda.fill.fore_color.rgb = _AZUL
            celda.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            celda.text_frame.paragraphs[0].font.size = Pt(11)
            celda.text_frame.paragraphs[0].font.bold = True
        for i, fila in enumerate(trozo, start=1):
            for j in range(len(columnas)):
                v = fila[j] if j < len(fila) else ""
                if isinstance(v, float):
                    v = f"{v:,.2f}"
                c = tabla.cell(i, j)
                c.text = "" if v is None else str(v)
                c.text_frame.paragraphs[0].font.size = Pt(10)


def _lamina_cierre(prs: Presentation, sql: str) -> None:
    lam = prs.slides.add_slide(prs.slide_layouts[6])
    linea = lam.shapes.add_shape(1, 0, Inches(0.2), prs.slide_width, Inches(0.12))
    linea.fill.solid()
    linea.fill.fore_color.rgb = _AMARILLO
    linea.line.fill.background()
    cuadro = lam.shapes.add_textbox(Inches(0.5), Inches(0.6), prs.slide_width - Inches(1.0), Inches(6.0))
    t = cuadro.text_frame
    t.word_wrap = True
    p = t.paragraphs[0]
    p.text = "SQL ejecutada (trazabilidad)"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = _AZUL
    p2 = t.add_paragraph()
    p2.text = sql[:1600]
    p2.font.size = Pt(9)
    p3 = t.add_paragraph()
    p3.text = ADVERTENCIA_LEGAL
    p3.font.size = Pt(9)
    p3.font.italic = True


def construir(pregunta: str, texto: str, sql: str, columnas: list[str], filas: list[list[Any]]) -> bytes:
    """Construye el .pptx (portada, tabla paginada y cierre con SQL + advertencia)."""
    prs = Presentation()
    _lamina_titulo(prs, pregunta, texto)
    if columnas:
        _laminas_tabla(prs, columnas, filas)
    _lamina_cierre(prs, sql)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
